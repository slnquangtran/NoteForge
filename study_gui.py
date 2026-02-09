import customtkinter as ctk
from tkinter import filedialog, messagebox
import threading
import os

# Import Logic
try:
    from study_assistant import LectureNoteGenerator
except ImportError:
    # Handle folder name conflict by importing the file directly
    import importlib.util
    import os
    if os.path.exists("study_assistant.py"):
        spec = importlib.util.spec_from_file_location("study_assistant_mod", "study_assistant.py")
        mod = importlib.util.module_from_spec(spec)
        spec.loader.exec_module(mod)
        LectureNoteGenerator = mod.LectureNoteGenerator
    else:
        LectureNoteGenerator = None

class StudyAssistantGUI(ctk.CTkToplevel):
    def __init__(self, master=None):
        super().__init__(master)
        
        self.title("NoteForge - Note Summarization")
        self.configure(fg_color="#312C51")
        self.center_window()
        
        self.assistant = None
        self.processed_data = None
        
        # Initialize assistant immediately (no heavy model loading)
        if LectureNoteGenerator:
            self.assistant = LectureNoteGenerator()
            status_text = "Ready. Load a file to generate lecture notes."
        else:
            status_text = "Error: missing logic module."
        
        self.status_var = ctk.StringVar(value=status_text)
        self.create_widgets()

    def center_window(self):
        self.update_idletasks()
        width = 1000
        height = 800
        x = (self.winfo_screenwidth() // 2) - (width // 2)
        y = (self.winfo_screenheight() // 2) - (height // 2)
        self.geometry(f"{width}x{height}+{x}+{y}")

    def create_widgets(self):
        self.grid_columnconfigure(0, weight=3) # Notes box (Left/Top)
        self.grid_rowconfigure(1, weight=3)    # Notes box weight
        self.grid_rowconfigure(2, weight=1)    # Thoughts box weight (Bottom)

        # --- Header Card ---
        header_card = ctk.CTkFrame(self, fg_color="#48426D", corner_radius=25)
        header_card.grid(row=0, column=0, padx=20, pady=(20, 10), sticky="ew")
        
        header = ctk.CTkFrame(header_card, fg_color="transparent")
        header.pack(fill="x", padx=20, pady=15)
        header.grid_columnconfigure(0, weight=1) # Title label expands
        header.grid_columnconfigure(1, weight=0) # Button doesn't expand
        
        ctk.CTkLabel(header, text="📝 AI Study Assistant", font=("Segoe UI", 20, "bold")).grid(row=0, column=0, padx=0, pady=0, sticky="w")
        
        self.btn_load = ctk.CTkButton(
            header, 
            text="📂 Load File (TXT/PPTX)", 
            command=self.load_file,
            font=("Segoe UI", 14),
            corner_radius=8,
            fg_color="#48426D", # Secondary Slate
            hover_color="#5A547F"
        )
        self.btn_load.grid(row=0, column=1, padx=0, pady=0, sticky="e")

        # --- Main Card ---
        main_card = ctk.CTkFrame(self, fg_color="#48426D", corner_radius=30)
        main_card.grid(row=1, column=0, padx=20, pady=5, sticky="nsew")

        self.notes_box = ctk.CTkTextbox(
            main_card, 
            font=("Segoe UI", 14), # Slightly larger
            wrap="word",
            corner_radius=20,
            border_width=2,
            border_color="#312C51"
        )
        self.notes_box.pack(fill="both", expand=True, padx=20, pady=20)

        # --- Thoughts Card ---
        thoughts_card = ctk.CTkFrame(self, fg_color="#48426D", corner_radius=25)
        thoughts_card.grid(row=3, column=0, padx=20, pady=(5, 10), sticky="nsew")

        ctk.CTkLabel(thoughts_card, text="🧠 AI Thought Process", font=("Segoe UI", 12, "italic"), text_color="#F0C38E").pack(padx=20, pady=(10, 0), anchor="w")
        
        self.thoughts_box = ctk.CTkTextbox(
            thoughts_card,
            height=120,
            font=("Segoe UI", 12),
            wrap="word",
            corner_radius=15,
            border_width=1,
            border_color="#312C51",
            fg_color="#1E1E1E",
            text_color="#A0A0A0"
        )
        self.thoughts_box.pack(fill="both", expand=True, padx=20, pady=(5, 15))
        self.thoughts_box.insert("1.0", "AI thoughts will appear here during processing...")
        self.thoughts_box.configure(state="disabled")

        # --- Footer Card ---
        footer_card = ctk.CTkFrame(self, fg_color="#48426D", corner_radius=25)
        footer_card.grid(row=4, column=0, padx=20, pady=(10, 20), sticky="ew")

        footer = ctk.CTkFrame(footer_card, fg_color="transparent")
        footer.pack(fill="x", padx=20, pady=15)
        footer.grid_columnconfigure(0, weight=0) # Progress bar
        footer.grid_columnconfigure(1, weight=1) # Status label (expands)
        footer.grid_columnconfigure(2, weight=0) # Process button
        footer.grid_columnconfigure(3, weight=0) # Export button
        
        self.progress_bar = ctk.CTkProgressBar(footer, width=200, height=10, corner_radius=5, fg_color="#48426D", progress_color="#F0C38E")
        self.progress_bar.grid(row=0, column=0, padx=(0, 15), pady=0, sticky="w")
        self.progress_bar.set(0)
        
        self.lbl_status = ctk.CTkLabel(footer, textvariable=self.status_var, font=("Segoe UI", 12), text_color="gray")
        self.lbl_status.grid(row=0, column=1, padx=0, pady=0, sticky="w")
        
        self.btn_export = ctk.CTkButton( # Moved export before process for typical flow
            footer, 
            text="💾 Export PDF", 
            command=self.export_pdf, 
            state="disabled",
            font=("Segoe UI", 14),
            corner_radius=8,
            fg_color="#F0C38E", # Tan/Peach
            hover_color="#DEB17E",
            text_color="#312C51"
        )
        self.btn_export.grid(row=0, column=2, padx=(0, 15), pady=0, sticky="e") # Pad to the left
        
        self.btn_process = ctk.CTkButton(
            footer, 
            text="📝 Generate Lecture Notes", 
            command=self.start_processing, 
            fg_color="#F1AA9B", # Coral
            hover_color="#DD998A",
            text_color="#312C51",
            font=("Segoe UI", 14),
            corner_radius=8
        )
        self.btn_process.grid(row=0, column=3, padx=0, pady=0, sticky="e")

        self.loaded_filepath = None

    def load_file(self):
        filepath = filedialog.askopenfilename(filetypes=[("Text Files", "*.txt"), ("PowerPoint", "*.pptx"), ("All Files", "*.*")])
        if filepath:
            self.loaded_filepath = filepath # Store path, don't read text yet if binary
            
            try:
                # Robust preview logic
                self.notes_box.delete("1.0", "end")
                self.thoughts_box.configure(state="normal")
                self.thoughts_box.delete("1.0", "end")
                
                if filepath.lower().endswith(".txt"):
                    # Try multiple encodings
                    content = ""
                    for enc in ['utf-8', 'latin-1', 'cp1252']:
                        try:
                            with open(filepath, "r", encoding=enc) as f:
                                content = f.read(2000) # Read more for preview
                            break
                        except: continue
                    
                    self.notes_box.insert("1.0", f"Preview (Text File):\n{'-'*40}\n{content}...")
                    self.thoughts_box.insert("1.0", f"Selected text file: {os.path.basename(filepath)}")
                
                elif filepath.lower().endswith(".pptx"):
                    from pptx import Presentation
                    prs = Presentation(filepath)
                    slide_count = len(prs.slides)
                    titles = []
                    for i, slide in enumerate(prs.slides[:5]):
                        title = slide.shapes.title.text if slide.shapes.title else f"Slide {i+1}"
                        titles.append(f"- {title}")
                    
                    summary = f"PowerPoint Presentation:\n- Total Slides: {slide_count}\n- First few slide titles:\n" + "\n".join(titles)
                    self.notes_box.insert("1.0", f"Preview (PPTX):\n{'-'*40}\n{summary}\n\n[Click 'Generate' to analyze all slides]")
                    self.thoughts_box.insert("1.0", f"Selected PowerPoint: {os.path.basename(filepath)} with {slide_count} slides.")
                
                self.thoughts_box.configure(state="disabled")
                self.status_var.set(f"Ready: {os.path.basename(filepath)}")
            except Exception as e:
                messagebox.showerror("Error", f"Failed to preview file: {e}")

    def start_processing(self):
        if not self.loaded_filepath: 
            return

        self.status_var.set("AI is starting to analyze your content...")
        self.progress_bar.set(0)
        self.btn_process.configure(state="disabled")
        
        # Reset thoughts box
        self.thoughts_box.configure(state="normal")
        self.thoughts_box.delete("1.0", "end")
        self.thoughts_box.insert("end", ">>> AI Thought Process Started <<<\n")
        self.thoughts_box.configure(state="disabled")
        
        threading.Thread(target=self.run_pipeline, daemon=True).start()

    def run_pipeline(self):
        if not self.assistant: return
        
        def update_status(msg, progress=0):
            self.status_var.set(msg)
            self.progress_bar.set(progress)
            
            # If it's a "thought", add it to the thoughts box
            if msg.startswith("AI Thought:"):
                self.after(0, lambda m=msg: self.add_thought(m))

        try:
            # Use process_file which handles parsing
            results = self.assistant.process_file(self.loaded_filepath, progress_callback=update_status)
            self.processed_data = results
            self.after(0, self.display_results)
        except Exception as e:
            import traceback
            error_details = traceback.format_exc()
            self.after(0, lambda: self.status_var.set(f"Error: {e}"))
            self.after(0, lambda: self.add_thought(f"ERROR: {e}\n{error_details}"))
            print(f"Pipeline Error:\n{error_details}")
            self.after(0, lambda: self.btn_process.configure(state="normal"))

    def add_thought(self, message):
        """Thread-safe update for the AI thoughts box"""
        def update():
            self.thoughts_box.configure(state="normal")
            # Strip "AI Thought: " for cleaner logging
            clean_msg = message.replace("AI Thought: ", "")
            self.thoughts_box.insert("end", f"→ {clean_msg}\n")
            self.thoughts_box.see("end") # Scroll to bottom
            self.thoughts_box.configure(state="disabled")
        
        self.after(0, update)

    def display_results(self):
        data = self.processed_data
        
        if not data or "error" in data:
            error_msg = data.get("error", "Unknown error") if data else "No data"
            messagebox.showerror("Processing Error", error_msg)
            self.btn_process.configure(state="normal")
            self.status_var.set("Error occurred.")
            return

        # Display the formatted lecture notes
        self.notes_box.delete("1.0", "end")
        self.notes_box.insert("1.0", data.get('notes', 'No notes generated'))
        
        self.btn_export.configure(state="normal")
        self.btn_process.configure(state="normal")
        self.status_var.set("Lecture notes generated successfully!")

    def export_pdf(self):
        if not self.processed_data: return
        path = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF", "*.pdf")])
        if path:
            self.assistant.export_to_pdf(self.processed_data, path)
            messagebox.showinfo("Export", f"Saved to {path}")

if __name__ == "__main__":
    app = StudyAssistantGUI()
    app.mainloop()
