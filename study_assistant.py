import os
import re
from datetime import datetime
from collections import defaultdict
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_LEFT
import nltk
from transformers import pipeline, BartForConditionalGeneration, BartTokenizer
import torch

# Download required NLTK data
try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt', quiet=True)

class LectureNoteGenerator:
    def __init__(self):
        self.file_path = None
        self.file_type = None
        
        # Legal topic keywords for detection
        self.topic_keywords = {
            'Offer and Acceptance': ['offer', 'acceptance', 'invitation to treat', 'postal rule', 'unilateral contract'],
            'Consideration': ['consideration', 'benefit', 'detriment', 'sufficient', 'adequate', 'past consideration'],
            'Promissory Estoppel': ['estoppel', 'promissory estoppel', 'shield not sword', 'equitable'],
            'Intention to Create Legal Relations': ['intention', 'legal relations', 'domestic', 'commercial'],
            'Certainty of Terms': ['certainty', 'vague', 'uncertain', 'agreement to agree'],
            'Capacity': ['capacity', 'minor', 'mental incapacity', 'intoxication'],
            'Privity of Contract': ['privity', 'third party', 'rights of third parties']
        }
        
        self.bart_summarizer = None  # Lazy loading
        self.bart_model_loaded = False
    
    def process_file(self, file_path, progress_callback=None):
        """Main entry point for generating lecture notes"""
        self.file_path = file_path
        self.file_type = self._detect_file_type()
        
        if progress_callback:
            progress_callback("Loading file...", 0.1)
        
        # Extract raw content
        if self.file_type == 'powerpoint':
            raw_content = self._extract_powerpoint_content()
        else:
            raw_content = self._extract_text_content()
        
        # Organize by topics
        if progress_callback:
            progress_callback("AI Thought: Organizing content by legal topics...", 0.4)
        topics = self._organize_by_topics(raw_content, progress_callback)
        
        if progress_callback:
            progress_callback("AI Thought: Generating AI-powered summaries with BART...", 0.6)

        # Generate BART summaries for topics
        self._enhance_topics_with_bart(topics, progress_callback)
        
        if progress_callback:
            progress_callback("Formatting notes...", 0.8)
        
        # Format as study guide
        formatted_notes = self._format_as_study_guide(topics, raw_content)
        
        if progress_callback:
            progress_callback("Complete!", 1.0)
        
        return {
            'notes': formatted_notes,
            'topics': topics,
            'ai_summaries': True,
            'bart_model_used': self.bart_model_loaded
        }

    def _initialize_bart_summarizer(self):
        """Lazy load BART model directly to avoid pipeline issues"""
        if not self.bart_model_loaded:
            try:
                model_name = "facebook/bart-large-cnn"
                self.tokenizer = BartTokenizer.from_pretrained(model_name)
                self.model = BartForConditionalGeneration.from_pretrained(model_name)
                
                # Multi-platform device detection
                if torch.cuda.is_available():
                    self.device = "cuda"
                    self.model.to(self.device)
                elif hasattr(torch.backends, "mps") and torch.backends.mps.is_available():
                    self.device = "mps"
                    self.model.to(self.device)
                else:
                    self.device = "cpu"
                
                print(f"BART using device: {self.device}")
                self.bart_model_loaded = True
            except Exception as e:
                import traceback
                error_msg = traceback.format_exc()
                print(f"Warning: Could not load BART model. Error:\n{error_msg}")
                self.model = None
                self.tokenizer = None

    def _clean_transcript_text(self, text):
        """Aggressively remove filler words, conversational markers, symbols, and parenthetical remarks"""
        if not text: return ""
        
        # 1. Remove parenthetical remarks like (oh yes, you)
        cleaned = re.sub(r'\([^)]*\)', '', text)
        
        # 2. Remove weird symbols and encoding artifacts
        cleaned = re.sub(r'[òùóéíáñçã]|\u00f2|\u00f9', '', cleaned)
        
        # 3. Define aggressive filler and conversational patterns
        fillers = [
            r'\b(um|uh|uhh|ah|err|hm|hmm)\b', 
            r'\b(you know|right|so basically|like|I mean|actually|obviously|literally|basically|essentially)\b',
            r'\b(I think|I\'m talking about|let\'s look at|we have|as I said|I want to say|believe me|I guess|you see)\b',
            r'\b(yes|no|okay|ok|actually|alright|yeah|yours|guess what)\b[.,!?]?',
            r'\b(does that make sense|is that clear|you see|you understand|look at it this way)\b[.,!?]?',
            r'\b(I am talking about|we are looking at|let me explain|let\'s move on|I should say)\b',
            r'\b(shame|wow|unreal|oh|uh-oh)\b[.,!?]?',
            r'\?+', # Remove question marks and entire questions if they seem rhetorical
        ]
        
        for pattern in fillers:
            cleaned = re.sub(pattern, '', cleaned, flags=re.IGNORECASE)
            
        # 4. Specifically target conversational questions
        sentences = nltk.sent_tokenize(cleaned)
        cleaned_sentences = []
        for s in sentences:
            s_strip = s.strip()
            # Filter out questions and very short fragments
            if not s_strip.endswith('?') and len(s_strip.split()) > 3:
                # Remove leading formal/conversational artifacts
                s_clean = re.sub(r'^(So,|Well,|And,|But,|Basically,|Actually,|Anyway,)\s+', '', s_strip, flags=re.IGNORECASE)
                cleaned_sentences.append(s_clean)
        
        cleaned = " ".join(cleaned_sentences)
        
        # 5. Clean up multiple spaces and punctuation artifacts
        cleaned = re.sub(r'\s+', ' ', cleaned)
        cleaned = re.sub(r'\s+([.,!?])', r'\1', cleaned)
        
        # 6. Final pass for starting punctuation
        cleaned = re.sub(r'^[.,!?\s]+', '', cleaned)
        
        return cleaned.strip()

    def _format_summary(self, summary_text):
        """Format a summary string into a clear bulleted list"""
        if not summary_text: return "No summary available."
        
        # Split into sentences using NLTK or regex
        sentences = nltk.sent_tokenize(summary_text)
        
        # Clean and format each sentence
        bullets = []
        for sent in sentences:
            sent = sent.strip()
            if len(sent) > 10:
                # Remove common conversational starts from summary too
                sent = re.sub(r'^(the speaker says|the lecture covers|this section is about)\s+', '', sent, flags=re.IGNORECASE)
                sent = sent[0].upper() + sent[1:] # Capitalize
                bullets.append(f"• {sent}")
        
        if not bullets:
            return summary_text
            
        return "\n".join(bullets)

    def _summarize_with_bart(self, text, max_length=150, min_length=50):
        """Use BART model directly to generate summary with cleaning and formatting"""
        if not text:
            return ""
        
        # Initialize BART if needed
        if not self.bart_model_loaded:
            self._initialize_bart_summarizer()
        
        # Helper to run inference on cleaned text
        def run_inference(input_texts):
            if not self.model or not self.tokenizer:
                return None
            try:
                inputs = self.tokenizer(input_texts, max_length=1024, return_tensors="pt", truncation=True, padding=True)
                input_ids = inputs["input_ids"]
                if self.device >= 0:
                    input_ids = input_ids.to(f"cuda:{self.device}")
                
                summary_ids = self.model.generate(input_ids, num_beams=4, max_length=max_length, min_length=min_length, early_stopping=True)
                return [self.tokenizer.decode(g, skip_special_tokens=True, clean_up_tokenization_spaces=False) for g in summary_ids]
            except Exception as e:
                print(f"BART Inference Error: {e}")
                return None

        # Handle Batch
        if isinstance(text, list):
            cleaned_texts = [self._clean_transcript_text(t) for t in text]
            summaries = run_inference(cleaned_texts)
            if summaries:
                return [self._format_summary(s) for s in summaries]
            return [self._format_summary(t[:300] + "...") for t in cleaned_texts]

        # Single string handling
        cleaned_text = self._clean_transcript_text(text)
        summaries = run_inference([cleaned_text])
        if summaries:
            return self._format_summary(summaries[0])
        return self._format_summary(cleaned_text[:300] + "...")
    def _enhance_topics_with_bart(self, topics, progress_callback=None):
        """Add BART-generated summaries to each topic using batching"""
        texts_to_summarize = []
        target_topics = []

        # First pass: identify what needs summarization
        for topic_name, content in topics.items():
            all_content = []
            for key in ['definitions', 'rules', 'cases', 'examples', 'exceptions']:
                all_content.extend(content[key])
            
            topic_text = " ".join(all_content)
            
            if len(topic_text) > 200:
                texts_to_summarize.append(topic_text)
                target_topics.append(topic_name)
            elif len(topic_text) > 0:
                print(f"DEBUG: Topic '{topic_name}' short, cleaning for summary.")
                content['bart_summary'] = self._format_summary(self._clean_transcript_text(topic_text))
            else:
                content['bart_summary'] = "No content available for summary."

        if not texts_to_summarize:
            return

        if progress_callback:
            progress_callback(f"AI Thought: Batch processing {len(texts_to_summarize)} summaries...", 0.65)

        # Batch summarize
        summaries = self._summarize_with_bart(texts_to_summarize)
        
        # If it returned a list, map them back
        if isinstance(summaries, list):
            for i, summary in enumerate(summaries):
                topic_name = target_topics[i]
                topics[topic_name]['bart_summary'] = summary
        else:
            # Fallback if it returned single string (unexpected)
            for topic_name in target_topics:
                topics[topic_name]['bart_summary'] = summaries

    def _generate_executive_summary(self, raw_content):
        """Generate overall executive summary using BART"""
        # Use first 2000 characters for executive summary
        excerpt = raw_content[:2000]
        return self._summarize_with_bart(excerpt, max_length=200, min_length=100)
    
    def _detect_file_type(self):
        """Detect if file is PowerPoint or text"""
        return 'powerpoint' if self.file_path.endswith('.pptx') else 'text'
    
    def _extract_powerpoint_content(self):
        """Extract all text from PowerPoint"""
        prs = Presentation(self.file_path)
        all_text = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    all_text.append(shape.text.strip())
            
            # Include speaker notes
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    all_text.append(notes_slide.notes_text_frame.text.strip())
        
        return ' '.join(all_text)
    
    def _extract_text_content(self):
        """Extract and clean text content"""
        with open(self.file_path, 'r', encoding='utf-8') as f:
            text = f.read()
        
        # Clean transcript
        return self._clean_transcript(text)
    
    def _clean_transcript(self, text):
        """Aggressively clean conversational transcript"""
        # Remove filler words and phrases
        fillers = [
            r'\bum+\b', r'\buh+\b', r'\bah+\b', r'\ber+\b', r'\blike\b',
            r'\byou know\b', r'\bI mean\b', r'\bbasically\b', r'\bactually\b',
            r'\bliterally\b', r'\bkind of\b', r'\bsort of\b', r'\bI think\b',
            r'\bI guess\b', r'\bI believe\b', r'\bI feel like\b', r'\bso yeah\b',
            r'\bokay so\b', r'\balright so\b', r'\bwell\b', r'\bso\b(?=\s+\w)',
            r'\bjust\b', r'\breally\b', r'\bvery\b', r'\bquite\b',
            r'\blet me\b', r'\blet\'s\b', r'\bgoing to\b', r'\bgonna\b',
            r'\bOkay,?\s*', r'\bAlright,?\s*', r'\bNow,?\s*'
        ]
        
        for filler in fillers:
            text = re.sub(filler, '', text, flags=re.IGNORECASE)
        
        # Remove repeated words
        text = re.sub(r'\b(\w+)\s+\1\b', r'\1', text, flags=re.IGNORECASE)
        
        # Clean up whitespace
        text = re.sub(r'\s+', ' ', text)
        text = re.sub(r'\s+([.,!?])', r'\1', text)
        
        return text.strip()
    
    def _organize_by_topics(self, text, progress_callback=None):
        """Organize content by legal topics"""
        # Split into sentences
        sentences = re.split(r'[.!?]+', text)
        sentences = [s.strip() for s in sentences if len(s.strip()) > 20]
        
        if progress_callback:
            progress_callback(f"AI Thought: Identified {len(sentences)} conceptual units for analysis.", 0.45)
        
        # Initialize topic structure
        topics = {}
        for topic_name in self.topic_keywords.keys():
            topics[topic_name] = {
                'definitions': [],
                'rules': [],
                'cases': [],
                'examples': [],
                'exceptions': []
            }
        
        # Classify each sentence
        for i, sentence in enumerate(sentences):
            if progress_callback and i % 10 == 0:
                progress_callback(f"AI Thought: Classifying content... ({i}/{len(sentences)})", 0.45 + (i/len(sentences))*0.15)
            
            # Detect which topic this sentence belongs to
            detected_topics = self._detect_topics(sentence)
            
            # Classify content type
            content_type = self._classify_content_type(sentence)
            
            # Add to appropriate topic(s)
            for topic in detected_topics:
                if topic in topics:
                    topics[topic][content_type].append(sentence)
        
        # Remove empty topics
        original_count = len(topics)
        topics = {k: v for k, v in topics.items() if any(v.values())}
        
        if progress_callback:
            progress_callback(f"AI Thought: Filtered {original_count - len(topics)} empty topics. Proceeding with {len(topics)} relevant legal areas.", 0.55)
        
        return topics
    
    def _detect_topics(self, sentence):
        """Detect which topics a sentence relates to"""
        sentence_lower = sentence.lower()
        detected = []
        
        for topic, keywords in self.topic_keywords.items():
            if any(keyword in sentence_lower for keyword in keywords):
                detected.append(topic)
        
        # If no topic detected, assign to first topic (general)
        return detected if detected else [list(self.topic_keywords.keys())[0]]
    
    def _classify_content_type(self, sentence):
        """Classify sentence as definition, rule, case, example, or exception"""
        sentence_lower = sentence.lower()
        
        # Check for legal cases
        if re.search(r'[A-Z][a-z]+\s+v\.?\s+[A-Z][a-z]+', sentence):
            return 'cases'
        
        # Check for definitions
        if any(word in sentence_lower for word in ['is defined as', 'means', 'refers to', 'is the', 'are called']):
            return 'definitions'
        
        # Check for exceptions
        if any(word in sentence_lower for word in ['exception', 'unless', 'however', 'but', 'except']):
            return 'exceptions'
        
        # Check for examples
        if any(word in sentence_lower for word in ['for example', 'e.g.', 'such as', 'for instance']):
            return 'examples'
        
        # Check for rules (must, should, requirement, rule)
        if any(word in sentence_lower for word in ['must', 'should', 'requirement', 'rule', 'principle']):
            return 'rules'
        
        # Default to rules
        return 'rules'
    
    def _format_as_study_guide(self, topics, raw_content):
        """Format organized topics as a study guide matching strict user requirements"""
        lines = []
        
        # Process each topic in strict order
        ordered_topics = [
            'Offer and Acceptance', 'Consideration', 'Promissory Estoppel',
            'Intention to Create Legal Relations', 'Certainty of Terms',
            'Capacity', 'Privity of Contract'
        ]
        
        for i, topic_key in enumerate(ordered_topics):
            # Check for the key regardless of case or slight variation
            actual_key = next((k for k in topics.keys() if k.lower() == topic_key.lower()), None)
            
            if not actual_key:
                continue
                
            content = topics[actual_key]
            
            # 1. Topic Header: ### [Number]. [TOPIC NAME IN ALL CAPS]
            lines.append(f"### {i+1}. {topic_key.upper()}")
            
            # 2. Overview: (1-2 clear sentences)
            summary = content.get('bart_summary', '')
            summary_sentences = [s for s in nltk.sent_tokenize(summary) if s.strip()]
            overview = " ".join(summary_sentences[:2]) if summary_sentences else "The legal principles governing this topic."
            lines.append(f"- **Overview:** {overview}")
            
            # Remaining summary becomes part of Rules or Other sections if needed
            additional_summary = " ".join(summary_sentences[2:])
            
            # 3. Key Definitions:
            lines.append("- **Key Definitions:**")
            if content['definitions']:
                for definition in content['definitions'][:5]:
                    lines.append(f"  - {definition}")
            else:
                lines.append("  - No specific definitions identified.")
            
            # 4. Legal Rules & Principles:
            lines.append("- **Legal Rules & Principles:**")
            rules_list = content['rules']
            if additional_summary:
                lines.append(f"  - {additional_summary}")
            
            if rules_list:
                for rule in rules_list[:10]:
                    cleaned_rule = rule.strip().lstrip('•').strip()
                    lines.append(f"  - {cleaned_rule}")
            elif not additional_summary:
                lines.append("  - Core legal principles as discussed in the lecture.")
            
            # 5. Exceptions & Special Cases:
            lines.append("- **Exceptions & Special Cases:**")
            if content['exceptions']:
                for exception in content['exceptions'][:5]:
                    cleaned_exc = exception.strip().lstrip('•').strip()
                    lines.append(f"  - {cleaned_exc}")
            else:
                lines.append("  - No significant exceptions noted in the transcript.")
            
            # 6. Practical Examples:
            lines.append("- **Practical Examples:**")
            if content['examples']:
                for example in content['examples'][:5]:
                    cleaned_ex = example.strip().lstrip('•').strip()
                    lines.append(f"  - {cleaned_ex}")
            else:
                lines.append("  - Concepts illustrated through general lecture discussion.")
            
            # 7. Key Cases/Statutes:
            lines.append("- **Key Cases/Statutes:**")
            if content['cases']:
                for case in content['cases'][:8]:
                    case_formatted = self._format_case_citation(case).strip().lstrip('•').strip()
                    lines.append(f"  - {case_formatted}")
            else:
                lines.append("  - General principles applied without specific case citation.")
            
            # 8. Study Tips:
            tips = self._generate_topic_tips(topic_key)
            lines.append(f"- **Study Tips:** {tips}")
            lines.append("")
        
        return '\n'.join(lines)

    def _generate_topic_tips(self, topic):
        """Provide specific study tips based on the legal topic"""
        tips = {
            'Offer and Acceptance': 'Always distinguish between an offer and an invitation to treat. Remember the postal rule requirements.',
            'Consideration': 'Always identify which promise is being enforced and what consideration was given for it. Use the benefit-detriment analysis.',
            'Promissory Estoppel': 'Remember: estoppel is a shield, not a sword. Focus on reliance and inequity.',
            'Intention to Create Legal Relations': 'Look at the relationship between parties and the context to determine intention.',
            'Certainty of Terms': 'Ask: "Is the term material? Can a court determine what was agreed?"',
            'Capacity': 'Capacity is rarely a central issue—but know the exceptions for minors and intoxicated persons.',
            'Privity of Contract': 'Remember the common law rule and the statutory exception for third parties.'
        }
        return tips.get(topic, "Focus on the core legal principle and key case law.")
    
    def _format_case_citation(self, text):
        """Format case citation with bold name"""
        # Extract case name
        match = re.search(r'([A-Z][a-z]+(?:\s+[A-Z][a-z]+)*)\s+v\.?\s+([A-Z][a-z]+(?:\s+[A-Z][a-z]+)*)', text)
        
        if match:
            case_name = f"{match.group(1)} v {match.group(2)}"
            
            # Extract year if present
            year_match = re.search(r'\((\d{4})\)', text)
            if year_match:
                case_name += f" ({year_match.group(1)})"
            
            # Replace in original text
            return text.replace(match.group(0), f"**{case_name}**")
        
        return text
    
    def export_to_pdf(self, data, filename):
        """Export lecture notes to PDF"""
        doc = SimpleDocTemplate(filename, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        # Parse the notes text
        notes_text = data.get('notes', '')
        lines = notes_text.split('\n')
        
        for line in lines:
            if line.startswith('TOPIC:'):
                story.append(Spacer(1, 12))
                story.append(Paragraph(line, styles['Heading1']))
            elif line.startswith(('DEFINITIONS:', 'LEGAL RULES:', 'KEY CASES:', 'EXCEPTIONS:', 'EXAMPLES:')):
                story.append(Paragraph(f"<b>{line}</b>", styles['Heading2']))
            elif line.strip():
                story.append(Paragraph(line, styles['Normal']))
                story.append(Spacer(1, 4))
        
        doc.build(story)

# For backward compatibility
ConceptualAssistant = LectureNoteGenerator
